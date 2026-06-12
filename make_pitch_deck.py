"""Generate the NestMatch AI investor pitch deck (.pptx).

Run:  python make_pitch_deck.py
Output: NestMatch_AI_Investor_Pitch.pptx

Content is grounded in the actual product. Bracketed [FILL IN] items are the
numbers only the founder can provide (market sizing, traction, financials, team)
— do not present fabricated figures. Every slide carries speaker notes.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Brand palette (matches the app) ----------------------------------------
PRIMARY = RGBColor(0x2F, 0x6B, 0xFF)   # indigo-blue
ACCENT = RGBColor(0x7C, 0x3A, 0xED)    # violet
TEAL = RGBColor(0x10, 0xB6, 0xA0)      # teal / "ai"
INK = RGBColor(0x0F, 0x17, 0x2A)       # near-black
MUTED = RGBColor(0x5B, 0x66, 0x7A)     # slate
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUCCESS = RGBColor(0x16, 0x9B, 0x62)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# --- Helpers ----------------------------------------------------------------
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    _solid(s, color)
    s.shadow.inherit = False
    return s


def rounded(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    _solid(s, color)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6):
    """runs: list of (string, size, color, bold) tuples = paragraphs."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, items, *, size=18, color=INK, gap=10):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (head, sub) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "▸  "
        r.font.size = Pt(size)
        r.font.color.rgb = TEAL
        r.font.name = FONT
        r.font.bold = True
        r2 = p.add_run()
        r2.text = head
        r2.font.size = Pt(size)
        r2.font.bold = True
        r2.font.color.rgb = color
        r2.font.name = FONT
        if sub:
            r3 = p.add_run()
            r3.text = "  —  " + sub
            r3.font.size = Pt(size - 3)
            r3.font.color.rgb = MUTED
            r3.font.name = FONT
    return box


def notes(slide, txt):
    slide.notes_slide.notes_text_frame.text = txt


def mark(slide, x, y, size=Inches(0.42)):
    """Small gradient-ish brand mark (rounded square)."""
    s = rounded(slide, x, y, size, size, PRIMARY)
    t = s.text_frame
    t.word_wrap = False
    p = t.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "◆"
    r.font.size = Pt(16)
    r.font.color.rgb = WHITE
    r.font.bold = True
    return s


def content_slide(title, kicker=None):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)
    rect(s, 0, 0, SW, Inches(0.16), PRIMARY)
    if kicker:
        text(s, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4),
             [(kicker.upper(), 12, TEAL, True)])
    text(s, Inches(0.7), Inches(0.72), Inches(12), Inches(0.9),
         [(title, 30, INK, True)])
    rounded(s, Inches(0.72), Inches(1.5), Inches(0.9), Inches(0.06), ACCENT)
    text(s, Inches(11.2), Inches(6.95), Inches(2), Inches(0.4),
         [("NestMatch AI", 9, MUTED, False)], align=PP_ALIGN.RIGHT)
    return s


def section_slide(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, INK)
    rect(s, 0, 0, Inches(0.28), SH, TEAL)
    rounded(s, Inches(5.0), Inches(2.0), Inches(3.3), Inches(0.5), ACCENT)
    text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2),
         [(title, 40, WHITE, True)])
    text(s, Inches(0.95), Inches(4.1), Inches(11), Inches(1),
         [(subtitle, 18, RGBColor(0xC8, 0xD0, 0xE0), False)])
    return s


# ============================================================================
# SLIDE 1 — TITLE
# ============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.28), SH, PRIMARY)
rect(s, Inches(0.28), 0, Inches(0.12), SH, ACCENT)
rect(s, Inches(0.40), 0, Inches(0.08), SH, TEAL)
mark(s, Inches(0.95), Inches(0.9), size=Inches(0.6))
text(s, Inches(1.7), Inches(0.92), Inches(8), Inches(0.6),
     [("NestMatch AI", 20, WHITE, True)])
text(s, Inches(0.95), Inches(2.7), Inches(11.5), Inches(1.6),
     [("The AI Operating System", 46, WHITE, True),
      ("for Real Estate", 46, TEAL, True)])
text(s, Inches(0.98), Inches(4.7), Inches(10.5), Inches(0.8),
     [("AI matchmaking · Off-plan marketplace · CRM · Branded proposals — one platform.",
       18, RGBColor(0xC8, 0xD0, 0xE0), False)])
rounded(s, Inches(0.95), Inches(5.7), Inches(4.6), Inches(0.7), PRIMARY)
text(s, Inches(0.95), Inches(5.78), Inches(4.6), Inches(0.55),
     [("Investment Pitch — Seeking [₹ AMOUNT]", 14, WHITE, True)], align=PP_ALIGN.CENTER)
text(s, Inches(0.98), Inches(6.7), Inches(11), Inches(0.4),
     [("[Founder Name] · [Role] · [Date] · [email/website]", 12, RGBColor(0x9A, 0xA6, 0xBC), False)])
notes(s, "Open in 15 seconds: 'We're building the AI operating system for real estate teams — "
         "CRM, off-plan inventory, and AI matchmaking in one platform, built to sell to 1,000 "
         "agencies and serve over a million buyers.' State the ask amount up front.")

# ============================================================================
# SLIDE 2 — PROBLEM
# ============================================================================
s = content_slide("Real estate teams are drowning in disconnected tools", kicker="The Problem")
bullets(s, Inches(0.75), Inches(1.9), Inches(11.8), Inches(4.5), [
    ("4–5 disconnected tools per agency", "CRM, portals, spreadsheets, WhatsApp, PDF makers — nothing talks to each other"),
    ("Leads leak and never get followed up", "[FILL IN: % of leads lost / hours wasted on manual follow-up]"),
    ("Off-plan information is fragmented", "Payment plans, availability, and developer data live in scattered brochures"),
    ("Buyers get generic, untrustworthy results", "Keyword search and gut-feel matching — no explanation of *why* a property fits"),
    ("AI is bolted on, not built in", "Most tools added a chatbot; none re-architected around intelligence"),
], size=19, gap=14)
notes(s, "Make the pain visceral and quantified. If you have a stat (e.g., 'agents spend X hours/week "
         "re-keying data' or '% of leads never contacted'), put it on the second bullet. Anchor on the "
         "agency's daily frustration, not abstract market problems.")

# ============================================================================
# SLIDE 3 — SOLUTION
# ============================================================================
s = content_slide("One AI-native platform that runs the whole brokerage", kicker="The Solution")
cards = [
    ("AI Matchmaking", "Hybrid recommender with explainable match reasons + DealScore™ investment grade", PRIMARY),
    ("Off-Plan Marketplace", "Developments, live unit availability & an AI payment-plan calculator", ACCENT),
    ("Smart CRM", "Lead capture, tagging, and auto-distribution (4 routing modes)", TEAL),
    ("Branded Proposals", "One-click PDF proposals in multiple languages & currencies", SUCCESS),
]
cx, cy, cw, ch = Inches(0.75), Inches(2.0), Inches(2.95), Inches(3.4)
for i, (t, d, c) in enumerate(cards):
    x = cx + i * (cw + Inches(0.12))
    rounded(s, x, cy, cw, ch, LIGHT)
    rounded(s, x, cy, cw, Inches(0.12), c)
    rounded(s, x + Inches(0.25), cy + Inches(0.35), Inches(0.5), Inches(0.5), c)
    text(s, x + Inches(0.25), cy + Inches(1.05), cw - Inches(0.5), Inches(0.7),
         [(t, 16, INK, True)])
    text(s, x + Inches(0.25), cy + Inches(1.7), cw - Inches(0.5), Inches(1.6),
         [(d, 12, MUTED, False)])
text(s, Inches(0.75), Inches(5.7), Inches(11.8), Inches(0.6),
     [("Powered by real Claude AI — with explainable outputs buyers and agents can trust.",
       16, INK, True)], align=PP_ALIGN.CENTER)
notes(s, "This is your one-sentence solution. Each card maps to a real, demoable feature. Land the line: "
         "'They digitized the workflow — we added the intelligence layer.'")

# ============================================================================
# SLIDE 4 — PRODUCT / WHAT'S BUILT
# ============================================================================
s = content_slide("A complete, production-grade platform — already built", kicker="Product")
col1 = [
    ("Multi-tenant B2B SaaS", "isolated data per agency; 58 production APIs"),
    ("Real Claude AI", "semantic search, copywriter, vision, assistant"),
    ("DealScore™", "one explainable 0–100 investment grade"),
    ("Off-plan + AI payment calculator", "booking → construction → handover"),
]
col2 = [
    ("CRM lead pipeline", "drag-and-drop stages + 4 distribution modes"),
    ("Natural-language search", "'3-bed under 900k with a pool' → results"),
    ("Security & compliance", "RBAC, API keys, audit logs, hashed creds"),
    ("Premium web app", "off-plan marketplace, dashboards, PDF proposals"),
]
bullets(s, Inches(0.75), Inches(2.0), Inches(6.0), Inches(4.2), col1, size=17, gap=16)
bullets(s, Inches(6.95), Inches(2.0), Inches(6.0), Inches(4.2), col2, size=17, gap=16)
notes(s, "Purpose of this slide: de-risk. 'This isn't a concept — it's built, tested (28 backend tests), "
         "and deployable today.' Technical investors love the 58 APIs / multi-tenant / audit-log facts.")

# ============================================================================
# SLIDE 5 — DEMO
# ============================================================================
s = section_slide("Live Demo", "Let me show you the product in 4 minutes.")
bullets(s, Inches(0.95), Inches(4.7), Inches(11.5), Inches(2.4), [
    ("Natural-language search", "type a sentence → parsed filters + ranked results"),
    ("Property page", "AI price intelligence, DealScore, one-click PDF proposal"),
    ("Off-plan project", "drag the AI payment-plan calculator live"),
    ("CRM pipeline", "switch to Fair Rotation, add a lead, watch it auto-assign"),
], size=15, color=RGBColor(0xC8, 0xD0, 0xE0), gap=8)
notes(s, "DEMO SCRIPT (rehearse cold, have a backup video): 1) seed data + set ANTHROPIC_API_KEY so AI is "
         "live. 2) NL search. 3) Property page: DealScore + generate PDF proposal. 4) Off-plan calculator "
         "(drag sliders). 5) CRM: Fair Rotation + add lead auto-assign + drag a card. Narrate VALUE, keep "
         "under 4 min. Optional: open /docs to show 58 APIs for technical investors.")

# ============================================================================
# SLIDE 6 — WHY NOW
# ============================================================================
s = content_slide("Why now", kicker="Timing")
bullets(s, Inches(0.75), Inches(2.0), Inches(11.8), Inches(4.2), [
    ("AI is finally production-ready", "LLMs (Claude) can do real extraction, vision & reasoning at scale — not just chat"),
    ("Off-plan boom", "UAE & India off-plan volumes surging; agencies need structured inventory + payment tooling"),
    ("Agencies are digitizing", "post-2024 shift from spreadsheets to SaaS across emerging real-estate markets"),
    ("Buyers expect intelligence", "consumers now demand explainable, personalized recommendations"),
], size=19, gap=18)
notes(s, "The 'why now' answers 'why didn't this exist before / why won't a giant just build it.' The honest "
         "answer: LLMs only recently became reliable + cheap enough to build the intelligence layer.")

# ============================================================================
# SLIDE 7 — MARKET
# ============================================================================
s = content_slide("A large, digitizing market", kicker="Market")
circles = [
    ("TAM", "[FILL IN]", "All real-estate brokerages globally × ARPU", PRIMARY, Inches(3.2)),
    ("SAM", "[FILL IN]", "UAE + India + North America beachhead", ACCENT, Inches(2.5)),
    ("SOM", "[FILL IN]", "Reachable in 3 years via founder-led + partner sales", TEAL, Inches(1.8)),
]
bx = Inches(1.1)
for label, val, desc, color, dia in circles:
    cyc = Inches(2.2) + (Inches(3.2) - dia) / 2
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, bx, cyc, dia, dia)
    _solid(c, color)
    c.shadow.inherit = False
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.bold = True; r.font.size = Pt(18); r.font.color.rgb = WHITE; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = val; r2.font.size = Pt(13); r2.font.color.rgb = WHITE; r2.font.name = FONT
    bx += dia + Inches(0.5)
text(s, Inches(7.2), Inches(2.4), Inches(5.6), Inches(3.5),
     [("TAM — [FILL IN: # brokerages × ARPU]", 16, INK, True),
      ("SAM — [FILL IN: serviceable markets]", 16, INK, True),
      ("SOM — [FILL IN: 3-yr reachable]", 16, INK, True),
      ("", 8, INK, False),
      ("Cite credible sources (industry reports, # of agencies per market, average SaaS ARPU). "
       "Bottom-up sizing (agencies × price) is more believable than top-down.", 12, MUTED, False)],
     space=12)
notes(s, "DO NOT fabricate. Build bottom-up: (# target agencies) × (your annual price) = SAM. Quote sources. "
         "Investors discount round top-down numbers; they trust a defensible bottom-up build.")

# ============================================================================
# SLIDE 8 — MOAT
# ============================================================================
s = content_slide("Defensibility — why we win and keep winning", kicker="Moat")
bullets(s, Inches(0.75), Inches(2.0), Inches(11.8), Inches(4.3), [
    ("Data flywheel", "every search, match & lead improves recommendations — compounding advantage"),
    ("Explainable AI", "DealScore + match reasons build buyer trust; hard to replicate credibly"),
    ("Off-plan depth", "structured developments + payment engine = high switching cost for agencies"),
    ("Multi-market by design", "currencies, languages, markets — not locked to one geography"),
    ("Integration lock-in", "portals, CRM, proposals in one workflow agencies run their business on"),
], size=18, gap=15)
notes(s, "Investors will ask 'what's the moat / why can't someone copy this.' Lead with the data flywheel "
         "(network effects within each tenant + across the platform). The model isn't the moat — the data and "
         "workflow lock-in are.")

# ============================================================================
# SLIDE 9 — BUSINESS MODEL
# ============================================================================
s = content_slide("Business model — B2B SaaS with AI usage metering", kicker="How we make money")
rows = [
    ("Plan", "Price/mo", "Seats", "Listings", "AI calls/mo"),
    ("Trial", "Free", "3", "50", "500"),
    ("Starter", "$199", "10", "1,000", "10,000"),
    ("Growth", "$799", "50", "25,000", "150,000"),
    ("Enterprise", "$2,499", "Unlimited", "Unlimited", "Unlimited"),
]
tbl_shape = s.shapes.add_table(len(rows), 5, Inches(0.75), Inches(2.0), Inches(8.4), Inches(3.2))
table = tbl_shape.table
for ci in range(5):
    table.columns[ci].width = Inches(1.68)
for ri, row in enumerate(rows):
    for cni, val in enumerate(row):
        cell = table.cell(ri, cni)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if cni else PP_ALIGN.LEFT
        run = para.runs[0]
        run.font.name = FONT
        run.font.size = Pt(13)
        if ri == 0:
            run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = INK
        else:
            run.font.bold = cni == 0
            run.font.color.rgb = INK
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
bullets(s, Inches(9.4), Inches(2.1), Inches(3.4), Inches(3.5), [
    ("Per-seat SaaS", "predictable recurring revenue"),
    ("AI usage metering", "expansion revenue as usage grows"),
    ("Enterprise + API", "high-ARPU agencies & partners"),
], size=14, gap=12)
text(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(1),
     [("[FILL IN: target ARPU, gross margin, CAC, LTV — even as estimates, label them clearly]",
       14, MUTED, True)])
notes(s, "These tiers/quotas are REAL (in the product's plan catalog). Add your unit economics: ARPU, gross "
         "margin (SaaS = high), estimated CAC and LTV. If estimates, say so. Highlight AI metering as a built-in "
         "expansion-revenue lever.")

# ============================================================================
# SLIDE 10 — TRACTION
# ============================================================================
s = content_slide("Traction & readiness", kicker="Where we are")
text(s, Inches(0.75), Inches(1.95), Inches(11.8), Inches(0.5),
     [("Product readiness (real, today):", 16, INK, True)])
bullets(s, Inches(0.9), Inches(2.5), Inches(11.6), Inches(1.8), [
    ("Production-grade platform built & tested", "multi-tenant backend, 58 APIs, 28 passing tests, premium web app"),
    ("Real AI integrated", "Claude-powered search, DealScore, copywriter, vision, assistant"),
    ("Deployable now", "Docker, CI, migrations, S3 media, Supabase/Postgres"),
], size=16, gap=12)
text(s, Inches(0.75), Inches(4.6), Inches(11.8), Inches(0.5),
     [("Market traction (your numbers):", 16, INK, True)])
bullets(s, Inches(0.9), Inches(5.15), Inches(11.6), Inches(1.6), [
    ("[FILL IN] design partners / pilot agencies", ""),
    ("[FILL IN] waitlist / LOIs / signed agencies", ""),
    ("[FILL IN] revenue or usage metrics, if any", ""),
], size=16, gap=10)
notes(s, "If pre-revenue, lead with PRODUCT READINESS as de-risking, then be honest about market traction "
         "(pilots, LOIs, waitlist). Never invent customers. 'Built and ready' is itself traction at this stage.")

# ============================================================================
# SLIDE 11 — COMPETITION
# ============================================================================
s = content_slide("How we compare", kicker="Competition")
rows = [
    ("Capability", "Incumbents / rem-app", "NestMatch AI"),
    ("Off-plan + payment plans", "Yes", "Yes + AI calculator"),
    ("CRM & lead distribution", "Yes", "Yes (4 modes)"),
    ("AI matchmaking", "Basic", "Hybrid + explainable"),
    ("DealScore investment grade", "No", "Yes"),
    ("Semantic / NL search", "No", "Yes (Claude)"),
    ("LLM assistant + AI copywriter", "No", "Yes"),
    ("Multi-market", "Single-market focus", "UAE + India + N. America"),
]
tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.6))
table = tbl_shape.table
table.columns[0].width = Inches(4.5)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(3.5)
for ri, row in enumerate(rows):
    for cni, val in enumerate(row):
        cell = table.cell(ri, cni)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT if cni == 0 else PP_ALIGN.CENTER
        run = para.runs[0]
        run.font.name = FONT
        run.font.size = Pt(13)
        if ri == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = INK
        else:
            run.font.bold = cni == 0
            cell.fill.solid()
            if cni == 2:
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF6, 0xF2)
                run.font.color.rgb = SUCCESS
                run.font.bold = True
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                run.font.color.rgb = INK
notes(s, "Be fair (don't trash competitors) but show clear differentiation: AI depth, explainability, "
         "DealScore, and multi-market. Line: 'They digitized the workflow; we added the intelligence layer — "
         "and we aren't locked to one geography.'")

# ============================================================================
# SLIDE 12 — GO TO MARKET
# ============================================================================
s = content_slide("Go-to-market", kicker="Distribution")
bullets(s, Inches(0.75), Inches(2.0), Inches(11.8), Inches(4.3), [
    ("Land agencies (B2B), users come free", "sell to the brokerage; their agents & buyers onboard bottom-up"),
    ("Founder-led sales first", "high-touch onboarding of design-partner agencies to nail the wedge"),
    ("Portal & ad integrations", "Bayut, Property Finder, Dubizzle, Google/Facebook Ads as acquisition + stickiness"),
    ("Partner / API channel", "white-label & API access for developers and large agencies"),
    ("Expand by market & tier", "beachhead market → adjacent geographies; Starter → Growth → Enterprise"),
], size=18, gap=15)
notes(s, "Show a credible first wedge: which market, which agency profile, how you reach the first 10–20. "
         "B2B land-and-expand with bottom-up user growth is the story. Name your first channel concretely.")

# ============================================================================
# SLIDE 13 — TEAM
# ============================================================================
s = content_slide("Team", kicker="Who we are")
slots = [("[Founder Name]", "[Role] — [1-line founder-market fit]"),
         ("[Co-founder / Key hire]", "[Role] — [credential]"),
         ("[Advisor]", "[Domain expertise]")]
x = Inches(0.85)
for name, role in slots:
    rounded(s, x, Inches(2.2), Inches(3.6), Inches(2.6), LIGHT)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), Inches(2.55), Inches(1.0), Inches(1.0))
    _solid(c, PRIMARY)
    c.shadow.inherit = False
    text(s, x + Inches(0.2), Inches(3.75), Inches(3.2), Inches(0.5),
         [(name, 16, INK, True)], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.2), Inches(4.2), Inches(3.2), Inches(0.6),
         [(role, 12, MUTED, False)], align=PP_ALIGN.CENTER)
    x += Inches(3.95)
text(s, Inches(0.75), Inches(5.3), Inches(11.8), Inches(0.8),
     [("[FILL IN: why this team will win — domain expertise, prior exits, technical depth, "
       "key hires the raise funds]", 14, MUTED, True)])
notes(s, "Investors back people. State founder-market fit crisply. If the team is thin today, name advisors "
         "and the exact roles the raise will hire. Don't leave this blank in the live pitch.")

# ============================================================================
# SLIDE 14 — THE ASK
# ============================================================================
s = content_slide("The ask", kicker="Investment")
rounded(s, Inches(0.75), Inches(1.95), Inches(5.7), Inches(1.5), PRIMARY)
text(s, Inches(1.0), Inches(2.15), Inches(5.2), Inches(1.1),
     [("Raising [₹ AMOUNT]", 26, WHITE, True),
      ("for [18-24] months of runway", 14, RGBColor(0xDD, 0xE6, 0xFF), False)])
text(s, Inches(0.75), Inches(3.7), Inches(6), Inches(0.5), [("Use of funds", 16, INK, True)])
bullets(s, Inches(0.9), Inches(4.2), Inches(6.0), Inches(2.4), [
    ("~40% Engineering & AI", ""),
    ("~35% Sales & GTM", ""),
    ("~15% Data & infrastructure", ""),
    ("~10% G&A", ""),
], size=15, gap=10)
text(s, Inches(7.1), Inches(3.7), Inches(5.5), Inches(0.5), [("Milestones this funds", 16, INK, True)])
bullets(s, Inches(7.25), Inches(4.2), Inches(5.4), Inches(2.4), [
    ("[FILL IN] paying agencies", ""),
    ("[FILL IN] ARR", ""),
    ("[FILL IN] markets live", ""),
    ("[FILL IN] AI accuracy / engagement target", ""),
], size=15, gap=10)
notes(s, "Right-size the ask to your stage and traction — a very large ask with no revenue invites pushback. "
         "Tie the amount to concrete milestones that set up the NEXT round at a higher valuation. End the pitch "
         "HERE (the ask), not on a thank-you slide.")

# ============================================================================
# SLIDE 15 — CLOSING
# ============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, 0, Inches(0.28), SH, TEAL)
mark(s, Inches(0.95), Inches(0.9), size=Inches(0.55))
text(s, Inches(1.65), Inches(0.92), Inches(8), Inches(0.6), [("NestMatch AI", 18, WHITE, True)])
text(s, Inches(0.95), Inches(2.9), Inches(11.5), Inches(1.6),
     [("The intelligence layer", 42, WHITE, True),
      ("for global real estate.", 42, TEAL, True)])
bullets(s, Inches(0.98), Inches(4.6), Inches(11), Inches(1.4), [
    ("Built, tested, and deployable today", ""),
    ("Real AI · multi-market · multi-tenant SaaS", ""),
], size=16, color=RGBColor(0xC8, 0xD0, 0xE0), gap=8)
text(s, Inches(0.98), Inches(6.4), Inches(11), Inches(0.6),
     [("[Founder Name] · [email] · [phone] · [website]", 14, RGBColor(0xC8, 0xD0, 0xE0), False)])
notes(s, "Closing: restate the one-liner and the ask, then go to Q&A. Have the competition matrix, financial "
         "model, and a recorded demo ready as backup. Confident, specific, numbers over adjectives.")

OUT = "NestMatch_AI_Investor_Pitch.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides.")
